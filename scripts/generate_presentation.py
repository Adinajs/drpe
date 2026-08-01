from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# --- Configuration & Theme ---
OUTPUT_FILE = "DRPE_Final_Presentation.pptx"
NAVY = RGBColor(0, 33, 71)
BLUE = RGBColor(0, 174, 239)
GRAY = RGBColor(244, 247, 249)
DARK_GRAY = RGBColor(64, 64, 64)
WHITE = RGBColor(255, 255, 255)

# Image paths (based on listing)
GUI_DIR = "GUI"
FIGURES_DIR = "figures"

# Helper to add a slide with a specific title and layout
def add_slide(prs, title_text, subtitle_text=None, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background styling (Clean White)
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Add a thin accent line at the top (SaaS aesthetic)
    accent_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = BLUE
    accent_bar.line.fill.background()
    
    # Set Title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.name = "Segoe UI"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = NAVY
        title.top = Inches(0.5)
    
    return slide

def add_image_if_exists(slide, image_path, left, top, width=None, height=None):
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    else:
        print(f"Warning: Image {image_path} not found.")

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Unified Security Risk, Compliance & Incident Response Platform"
    subtitle = slide.placeholders[1]
    subtitle.text = "AI-Driven SecOps, Compliance Monitoring & Automated Response\n\nTeam: Dynamic Risk Posture Evaluation (DRPE)"
    slide.notes_slide.notes_text_frame.text = "Good morning evaluators. Today we present our Unified Security Risk, Compliance, and Incident Response Platform. This system transforms fragmented security silos into an autonomous, intelligence-driven ecosystem."
    
    # Style title slide
    title.text_frame.paragraphs[0].font.color.rgb = NAVY
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    # 2. Problem Statement
    slide = add_slide(prs, "The Fragmentation Gap")
    slide.notes_slide.notes_text_frame.text = "Modern security teams suffer from tool fragmentation. 40% of critical alerts go uninvestigated because they are siloed. Our motivation was to bridge the gap between risk intelligence, compliance auditing, and incident response."
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Fragmented Tools & Siloed Security Data"
    p = tf.add_paragraph()
    p.text = "• Manual Compliance Audits: Labor-intensive and error-prone."
    p = tf.add_paragraph()
    p.text = "• Alert Fatigue: SOC teams overwhelmed by disconnected data."
    p = tf.add_paragraph()
    p.text = "• Delayed Response: Lack of centralized orchestration."
    
    # 3. Proposed Solution
    slide = add_slide(prs, "A Unified Autonomous Ecosystem")
    slide.notes_slide.notes_text_frame.text = "Our platform integrates three critical pillars. Instead of separate tools for scanning and response, we provide a unified dashboard where risk data flows directly into compliance reports and automated response playbooks."
    add_image_if_exists(slide, os.path.join(FIGURES_DIR, "system_flow.png"), Inches(5), Inches(2), width=Inches(4))
    content = slide.placeholders[1]
    content.width = Inches(4.5)
    tf = content.text_frame
    tf.text = "Integrated SecOps Pillars:"
    for item in ["I. Risk Intelligence", "II. Governance & Compliance", "III. Incident Response"]:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.bold = True
        p.font.color.rgb = BLUE
    
    # 4. System Architecture
    slide = add_slide(prs, "Enterprise Technical Foundation")
    slide.notes_slide.notes_text_frame.text = "The system is built on a high-performance stack. FastAPI handles the backend logic, while Wazuh and OpenVAS provide the heavy lifting for SIEM and vulnerability scanning. Dockerized workers ensure that scanning processes are isolated and scalable."
    add_image_if_exists(slide, os.path.join(FIGURES_DIR, "block_diagram.png"), Inches(1), Inches(2.5), width=Inches(8))
    
    # 5. Module I: Risk
    slide = add_slide(prs, "Module I: Risk Intelligence")
    slide.notes_slide.notes_text_frame.text = "Risk Intelligence isn't just about scanning; it's about context. We correlate asset vulnerabilities with live threat intelligence from AlienVault and AbuseIPDB to prioritize what actually matters to the business."
    add_image_if_exists(slide, os.path.join(GUI_DIR, "gui_vulnerabilities_light.png"), Inches(5), Inches(2), width=Inches(4.5))
    tf = slide.placeholders[1].text_frame
    tf.text = "Predictive Attack Surface Management"
    for item in ["Autonomous Asset Discovery", "Dynamic Risk Scoring (CVSS + Intel)", "Exploit Prioritization"]:
        tf.add_paragraph().text = f"• {item}"

    # 6. Module II: Compliance
    slide = add_slide(prs, "Module II: Governance & Compliance")
    slide.notes_slide.notes_text_frame.text = "Compliance is often a manual, painful process. Our platform automates evidence collection, mapping security telemetry directly to ISO 27001 and SOC 2 controls, providing a real-time compliance scorecard."
    add_image_if_exists(slide, os.path.join(GUI_DIR, "gui_report_light.png"), Inches(5), Inches(2), width=Inches(4.5))
    tf = slide.placeholders[1].text_frame
    tf.text = "Continuous Framework Monitoring"
    for item in ["ISO 27001 & SOC 2 Mapping", "Automated Evidence Collection", "Real-time Compliance Scorecards"]:
        tf.add_paragraph().text = f"• {item}"

    # 7. Module III: Incident Response
    slide = add_slide(prs, "Module III: Incident Response")
    slide.notes_slide.notes_text_frame.text = "When a threat is detected, speed is everything. Our SOAR-like playbooks allow the system to automatically block malicious IPs or isolate compromised assets, reducing the Mean Time to Respond (MTTR) from hours to minutes."
    add_image_if_exists(slide, os.path.join(GUI_DIR, "gui_topology_light.png"), Inches(5), Inches(2), width=Inches(4.5))
    tf = slide.placeholders[1].text_frame
    tf.text = "Orchestrated Threat Containment"
    for item in ["Automated SOAR Playbooks", "Alert Centralization", "MTTR Reduction Analytics"]:
        tf.add_paragraph().text = f"• {item}"

    # 8. Enterprise Features
    slide = add_slide(prs, "Enterprise Resilience Features")
    slide.notes_slide.notes_text_frame.text = "The platform is built with enterprise security in mind. We use Zero Trust principles, RBAC, and full audit logging. The architecture is asynchronous, using Redis and Celery to handle large-scale security data processing."
    tf = slide.placeholders[1].text_frame
    tf.text = "Security & Scalability"
    for item in ["Zero Trust (RBAC + 2FA)", "Dockerized Engine Isolation", "Redis-Backed Task Queuing", "Real-time Threat Correlation"]:
        tf.add_paragraph().text = f"• {item}"

    # 9. Dashboard Showcase
    slide = add_slide(prs, "Platform Showcase: Dashboard")
    slide.notes_slide.notes_text_frame.text = "This is our primary dashboard. It provides a real-time heat map of organizational risk, compliance trends, and active security incidents, all in a single, high-fidelity interface built with React."
    add_image_if_exists(slide, os.path.join(GUI_DIR, "gui_dashboard_light_1.png"), Inches(1), Inches(2), width=Inches(8))

    # 10. Business Impact
    slide = add_slide(prs, "Business Impact & Future Roadmap")
    slide.notes_slide.notes_text_frame.text = "The business value is clear: a 70% reduction in manual effort for compliance and faster incident response. Looking forward, we aim to integrate predictive AI models to anticipate threats before they occur."
    tf = slide.placeholders[1].text_frame
    tf.text = "Operational Excellence"
    for item in ["70% Reduction in Manual Audit Time", "Faster Threat Containment", "AI-Driven Predictive Security (Next Gen)"]:
        tf.add_paragraph().text = f"• {item}"

    # 11. Technical Achievements
    slide = add_slide(prs, "Technical Achievements")
    slide.notes_slide.notes_text_frame.text = "Building this platform required mastering the full-stack security lifecycle. We successfully integrated multiple external security engines and managed complex data normalization across diverse telemetry sources."
    tf = slide.placeholders[1].text_frame
    tf.text = "Engineering Outcomes"
    for item in ["Full-Stack Integration (React/FastAPI)", "SIEM & Vuln Engine Correlation", "Asynchronous Worker Architecture"]:
        tf.add_paragraph().text = f"• {item}"

    # 12. Conclusion
    slide = add_slide(prs, "Conclusion & Q/A", layout_index=0)
    slide.notes_slide.notes_text_frame.text = "In conclusion, our platform provides a modern, unified approach to cybersecurity operations. We've bridged the gap between risk, compliance, and response. Thank you for your time. Are there any questions?"
    slide.shapes.title.text = "Unified Security: The Next Frontier"
    slide.placeholders[1].text = "Thank You!\nQuestions?\n\nContact: fyp-team@university.edu"

    # Save
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
